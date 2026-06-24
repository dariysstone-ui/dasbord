#!/usr/bin/env python3
"""
Generate metrics PPTX — called from Node.js API with JSON data piped via stdin
Usage: python3 generate_pptx.py < data.json > output.pptx
"""
import sys, json, io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

def rgb(hex_str):
    return RGBColor(int(hex_str[0:2],16), int(hex_str[2:4],16), int(hex_str[4:6],16))

def set_cell_bg(cell, hex_color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # Remove existing fills
    for tag in ['a:solidFill','a:gradFill','a:noFill']:
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', hex_color)

def set_cell_border(cell, hex_color='BFC9CA', width_pt=0.5):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    w = int(width_pt * 12700)
    for side in ['a:lnL','a:lnR','a:lnT','a:lnB']:
        existing = tcPr.find(qn(side))
        if existing is not None:
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(side))
        ln.set('w', str(w))
        solidFill = etree.SubElement(ln, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', hex_color)

def set_cell_text(cell, text, font_name='Aptos Narrow', font_size_pt=11,
                  bold=False, color_hex='000000', align='center', valign='bottom',
                  margin_l_inch=0.05):
    tf = cell.text_frame
    tf.word_wrap = True
    # Clear existing
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ''
    para = tf.paragraphs[0]
    para.text = ''
    # Alignment
    algn_map = {'center': PP_ALIGN.CENTER, 'left': PP_ALIGN.LEFT, 'right': PP_ALIGN.RIGHT}
    para.alignment = algn_map.get(align, PP_ALIGN.CENTER)
    run = para.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(font_size_pt)
    font.bold = bold
    font.color.rgb = rgb(color_hex)
    # Vertical align
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcPr.set('anchor', 'b' if valign == 'bottom' else 'ctr')
    # Left margin for OMSU column
    tcPr.set('marL', str(int(margin_l_inch * 914400)))
    tcPr.set('marR', str(int(0.05 * 914400)))
    tcPr.set('marT', str(int(0.04 * 914400)))
    tcPr.set('marB', '0')

def build_table(slide, rows_data, x_in, y_in, col_widths_in, row_heights_in,
                period_str, HDR_COLOR='2E3F5C', HDR_TEXT='FFFFFF'):
    """Add a properly formatted table to slide."""
    from pptx.util import Inches
    total_w = sum(Inches(w) for w in col_widths_in)
    total_h = sum(Inches(h) for h in row_heights_in)

    tbl_placeholder = slide.shapes.add_table(
        len(row_heights_in), len(col_widths_in),
        Inches(x_in), Inches(y_in), total_w, total_h
    )
    tbl = tbl_placeholder.table

    # Set column widths
    for ci, w in enumerate(col_widths_in):
        tbl.columns[ci].width = Inches(w)
    # Set row heights
    for ri, h in enumerate(row_heights_in):
        tbl.rows[ri].height = Inches(h)

    # Header row
    HDR_LABELS = ['ОМСУ', f'Допустимое\nзначение\n{period_str}', f'Поступило\n{period_str}', 'Отклонение']
    for ci, lbl in enumerate(HDR_LABELS):
        cell = tbl.cell(0, ci)
        set_cell_bg(cell, HDR_COLOR)
        set_cell_border(cell, 'FFFFFF', 0.5)
        set_cell_text(cell, lbl, font_name='Arial', font_size_pt=11,
                      bold=True, color_hex=HDR_TEXT, align='center', valign='bottom',
                      margin_l_inch=0.04)

    # Data rows
    for ri, row in enumerate(rows_data):
        omsu, metric, actual, dev_pct = row
        bg = 'FFFFFF' if ri % 2 == 0 else 'F2F3F4'
        # Dev color
        if dev_pct > 20:   dev_color = 'C00000'
        elif dev_pct > 10: dev_color = 'C07000'
        else:              dev_color = '1E7B1E'
        dev_str = ('+' if dev_pct >= 0 else '') + str(dev_pct) + '%'

        cells_data = [
            (str(omsu),   'Aptos Narrow', False, '000000', 'left'),
            (str(metric), 'Aptos Narrow', False, '000000', 'center'),
            (str(actual), 'Aptos Narrow', False, '000000', 'center'),
            (dev_str,     'Aptos Narrow', True,  dev_color,'center'),
        ]
        for ci, (txt, font, bold, color, align) in enumerate(cells_data):
            cell = tbl.cell(ri + 1, ci)
            set_cell_bg(cell, bg)
            set_cell_border(cell, 'BFC9CA', 0.4)
            set_cell_text(cell, txt, font_name=font, font_size_pt=11,
                          bold=bold, color_hex=color, align=align, valign='bottom',
                          margin_l_inch=0.12 if ci == 0 else 0.04)

    return tbl_placeholder

def generate(data):
    period_str = data['period']
    title_str  = data['title']
    rows = data['rows']  # list of [omsu, metric, actual, dev_pct]

    # 16:9 slide
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.500)

    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)

    # ── Title block ──
    from pptx.shapes.autoshape import Shape as PShape
    from pptx.util import Inches, Pt
    title_h = Inches(0.775)
    # Blue rectangle
    txBox = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), Inches(13.333), title_h
    )
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = rgb('1B4F8A')
    txBox.line.fill.background()
    tf = txBox.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.text = title_str
    para.alignment = PP_ALIGN.LEFT
    run = para.runs[0]
    run.font.name = 'Arial'
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = rgb('FFFFFF')
    tc = txBox._element
    # vertical center
    txBody = tf._txBody
    txBody.set('anchor', 'ctr')
    # padding
    for attr, val in [('marL','228600'),('marR','114300'),('marT','45720'),('marB','45720')]:
        txBody.set(attr, val)

    # ── Split rows ──
    half = (len(rows) + 1) // 2
    left_rows  = rows[:half]
    right_rows = rows[half:]

    # Col widths matching reference
    col_w = [1.693, 1.515, 1.163, 1.445]
    hdr_h  = 0.716
    data_h = 0.224

    y_tbl = 0.775 + 0.05  # just below title

    left_rh  = [hdr_h] + [data_h] * len(left_rows)
    right_rh = [hdr_h] + [data_h] * len(right_rows)

    build_table(slide, left_rows,  0.10, y_tbl, col_w, left_rh,  period_str)
    build_table(slide, right_rows, 0.10 + sum(col_w) + 0.12, y_tbl, col_w, right_rh, period_str)

    # Slide number
    txb = slide.shapes.add_textbox(Inches(12.8), Inches(7.1), Inches(0.5), Inches(0.35))
    tf2 = txb.text_frame
    tf2.paragraphs[0].text = '1'
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
    run2 = tf2.paragraphs[0].runs[0]
    run2.font.size = Pt(10)
    run2.font.color.rgb = rgb('888888')

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

if __name__ == '__main__':
    data = json.loads(sys.stdin.read())
    result = generate(data)
    sys.stdout.buffer.write(result)
